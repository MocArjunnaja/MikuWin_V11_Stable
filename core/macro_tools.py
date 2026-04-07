import os
import json
import urllib.parse
from pathlib import Path
from core.tools import registry
import pyautogui
import time
import subprocess

try:
    from docx import Document
except ImportError:
    print("[MacroTools] Warning: python-docx not installed. Word macros will fail.")

try:
    from pptx import Presentation
except ImportError:
    print("[MacroTools] Warning: python-pptx not installed. PPT macros will fail.")

# Definisikan Output direktori agar file buatan AI tersimpan rapi
OUTPUT_DIR = Path("C:/Users/Arjun/Desktop/Miku_Documents")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

class MacroTools:
    """Kumpulan Makro (High-Level Actions) untuk MikuWin v10"""

    def __init__(self):
        pass

    @registry.register(name="create_word_document", description="Membantu user membuat file dokumen Microsoft Word (.docx) secara kilat tanpa harus buka aplikasinya.")
    def create_word_document(self, filename: str, content: str) -> str:
        try:
            if not filename.endswith('.docx'):
                filename += '.docx'
            
            doc = Document()
            # Pisahkan paragraf berdasarkan baris baru jika ada
            for paragraph in content.split('\n'):
                if paragraph.strip():
                    doc.add_paragraph(paragraph.strip())
            
            filepath = OUTPUT_DIR / filename
            doc.save(str(filepath))
            
            # Otomatis buka file jika ingin ditunjukkan ke user
            os.startfile(str(filepath))
            
            return f"Dokumen Word berhasil dibuat dan disimpan di {filepath}"
        except Exception as e:
            return f"Gagal membuat dokumen Word: {e}"

    @registry.register(name="read_word_document", description="Membaca isi file dokumen Microsoft Word (.docx) di Desktop agar AI dapat me-review atau mengecek isinya sebelum merevisi.")
    def read_word_document(self, filename: str) -> str:
        try:
            if not filename.endswith('.docx'):
                filename += '.docx'
            filepath = OUTPUT_DIR / filename
            if not filepath.exists():
                return f"File '{filename}' tidak ditemukan di {OUTPUT_DIR}."
            
            doc = Document(str(filepath))
            full_text = []
            for i, para in enumerate(doc.paragraphs):
                # Membantu AI mengetahui letak barisnya
                full_text.append(f"[Paragraf {i+1}]: {para.text}")
                
            return "\n".join(full_text)
        except Exception as e:
            return f"Gagal membaca dokumen Word: {e}"

    @registry.register(name="edit_word_document", description="Bertindak sebagai Microsoft Word Editor otonom (Agentic API). action bisa berupa: 'append' (tambahkan teks di akhir dokumen), 'replace' (tulis ulang / revisi teks yang ada), atau 'delete'. Untuk merevisi, target_text harus diisi tulisan lamanya, lalu new_text diisi tulisan baru. Parameter style opsional: 'bold', 'italic'.")
    def edit_word_document(self, filename: str, action: str, target_text: str = "", new_text: str = "", style: str = "") -> str:
        try:
            if not filename.endswith('.docx'):
                filename += '.docx'
            filepath = OUTPUT_DIR / filename
            if not filepath.exists():
                return f"File '{filename}' tidak ditemukan. Buat dulu dengan create_word_document."
            
            doc = Document(str(filepath))
            changes_made = False

            if action == "append":
                if new_text:
                    p = doc.add_paragraph()
                    run = p.add_run(new_text)
                    if style == "bold": run.bold = True
                    if style == "italic": run.italic = True
                    changes_made = True
            
            elif action in ["replace", "delete"]:
                if not target_text:
                    return "Gagal: Anda lupa menyertakan target_text yang ingin dihapus/direvisi."
                
                # Membaca seluruh paragraf
                for para in doc.paragraphs:
                    if target_text in para.text:
                        if action == "replace":
                            # Ubah isi paragraf
                            para.text = para.text.replace(target_text, new_text)
                            # Berikan styling pada hasil buatan baru (Simplifikasi ke seluruh paragraf)
                            if style == "bold":
                                for r in para.runs: r.bold = True
                            if style == "italic":
                                for r in para.runs: r.italic = True
                        elif action == "delete":
                            para.text = para.text.replace(target_text, "")
                        changes_made = True
            
            if changes_made:
                doc.save(str(filepath))
                # Buka file secara native di foreground agar user tahu sudah di-update
                os.startfile(str(filepath))
                return f"Berhasil mengedit dokumen '{filename}'. Teks sudah di-update di background."
            else:
                return f"Gagal: target_text '{target_text}' tidak ditemukan di dalam dokumen. Pastikan Anda sudah memanggil read_word_document sebelumnya."
        
        except Exception as e:
            return f"Gagal mengedit dokumen Word: {e}"

    @registry.register(name="create_powerpoint", description="Membantu user membuat file presentasi Microsoft PowerPoint (.pptx) sederhana.")
    def create_powerpoint(self, filename: str, title: str, content: str) -> str:
        try:
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title
            subtitle_shape.text = content
            
            filepath = OUTPUT_DIR / filename
            prs.save(str(filepath))
            
            os.startfile(str(filepath))
            
            return f"Presentasi PowerPoint berhasil dibuat dan disimpan di {filepath}"
        except Exception as e:
            return f"Gagal membuat presentasi PowerPoint: {e}"

    @registry.register(name="send_whatsapp_message", description="Mengirim pesan WhatsApp (Desktop) langsung ke nomor kontak yang tersimpan di address book tanpa harus mencari namanya. Parameter contact_name diisi dengan nama panggilan target.")
    def send_whatsapp_message(self, contact_name: str, message: str) -> str:
        try:
            print(f"[MacroTools] Mencoba mengirim WA ke '{contact_name}' ...")
            
            # Load Database Kontak (JSON)
            contacts_db_path = Path("data/contacts_db.json")
            if not contacts_db_path.exists():
                return f"Gagal: Database kontak (contacts_db.json) tidak ditemukan di folder data. Silakan parse CSV Anda dulu."
            
            with open(contacts_db_path, "r", encoding="utf-8") as f:
                contacts = json.load(f)
            
            # Cocokkan nama kontak (Case-insensitive)
            # Kita bisa pakai in untuk partial match (misal "Budi" cocok dengan "Budi Anto")
            target_phone = None
            search_name = contact_name.lower().strip()
            
            # 1. Exact Match dulu
            if search_name in contacts:
                target_phone = contacts[search_name]
            else:
                # 2. Partial Match (Pencarian nama yang mengandung kata tersebut)
                for db_name, phone in contacts.items():
                    if search_name in db_name:
                        target_phone = phone
                        print(f"[MacroTools] Ditemukan partial match: {db_name} -> {phone}")
                        contact_name = db_name # Perbarui ke nama lengkap yang ketemu
                        break
            
            if target_phone:
                # JALUR A: Kontak Personal (Menggunakan URI Scheme Canggih)
                print(f"[MacroTools] Nomor ditemukan: {target_phone}. Membuka WA Desktop...")
                text_encoded = urllib.parse.quote(message)
                
                uri = f"whatsapp://send?phone={target_phone}&text={text_encoded}"
                os.startfile(uri)
                
                time.sleep(4) # Tunggu WA Desktop selesai memuat draft
                pyautogui.press('enter')
                time.sleep(0.5)
                
                return f"Pesan WhatsApp berhasil dikirimkan langsung ke {contact_name} ({target_phone})."
            
            else:
                # JALUR B: Grup WA atau Kontak Belum Tersimpan (Fallback UI Automation)
                print(f"[MacroTools] '{contact_name}' tidak ada di buku telepon. Asumsi Grup WA. Beralih ke UI Automation...")
                
                # Buka WA Desktop saja (ke halaman utama)
                os.startfile("whatsapp://")
                time.sleep(3) # Tunggu animasi WA Windows selesai
                
                # Gunakan shortcut pencarian WA
                pyautogui.hotkey('ctrl', 'f')
                time.sleep(0.5)
                
                # Ketik nama grup secara perlahan
                pyautogui.write(contact_name, interval=0.05)
                time.sleep(2) # Tunggu hasil search WA merender
                
                # Tekan Enter untuk masuk ke room chat grup teratas
                pyautogui.press('enter')
                time.sleep(1)
                
                # Ketik pesannya
                pyautogui.write(message, interval=0.02)
                time.sleep(0.5)
                pyautogui.press('enter')
                
                return f"Pesan WhatsApp berhasil dikirimkan ke Grup/Non-Kontak '{contact_name}' menggunakan pencarian visual."
            
        except Exception as e:
            return f"Gagal mengirim pesan WA secara direct: {e}"

    @registry.register(name="send_telegram_message", description="Mengirim pesan Telegram desktop langsung ke seseorang menggunakan username (@username) atau nomor telepon.")
    def send_telegram_message(self, contact_name: str, message: str) -> str:
        try:
            print(f"[MacroTools] Mencoba mengirim Telegram ke '{contact_name}' ...")
            
            # Kita coba asumsikan contact_name adalah username jika diawali '@'
            # Jika tidak, idealnya kita cari di contacts_db.json juga jika mereka menggunakan nomor telepon yang sama.
            # Sayangnya Telegram Desktop URL Scheme (tg://) paling andal menggunakan `resolve?domain=` (untuk @username).
            
            target_username = contact_name.lower().strip()
            
            # Bersihkan tanda @ jika user tidak sengaja menyertakannya
            if target_username.startswith('@'):
                target_username = target_username[1:]
                
            # encode URL komponen pesan
            text_encoded = urllib.parse.quote(message)
            
            # Direct Injection ke Telegram Desktop
            # Format: tg://resolve?domain=USERNAME&text=PESAN
            uri = f"tg://resolve?domain={target_username}&text={text_encoded}"
            
            print(f"[MacroTools] Membuka URL Telegram: {uri}")
            os.startfile(uri)
            
            time.sleep(3) # Kasih waktu buka aplikasi dan loading chat room
            
            # Telegram, ketika dipanggil dengan tg://resolve?text= sudah menaruh text di draft chat.
            # Sama seperti WA, kita tinggal pencet enter
            pyautogui.press('enter')
            
            return f"Pesan Telegram berhasil dikirim ke @{target_username} langsung."
        except Exception as e:
            return f"Gagal mengirim pesan Telegram: {e}"
